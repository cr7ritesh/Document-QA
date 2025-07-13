import os
import logging
from flask import Flask, render_template, request, jsonify, session
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_cohere import CohereEmbeddings
from langchain.chat_models import init_chat_model
from langchain.chains import RetrievalQA
from dotenv import load_dotenv
import tempfile
from werkzeug.utils import secure_filename
import pytesseract
from PIL import Image
import io
import re
import requests
from bs4 import BeautifulSoup

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.DEBUG)

# Create Flask app
app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "fallback-secret-key")

# Configuration
UPLOAD_FOLDER = tempfile.mkdtemp()
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'png', 'jpg', 'jpeg'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_image_file(image_path):
    """Extract text from a standalone image file using OCR."""
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        logging.error(f"OCR error for image {image_path}: {str(e)}")
        return ""

def extract_text_from_image_bytes(image_bytes):
    """Extract text from image bytes (for embedded images)."""
    try:
        image = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        logging.error(f"OCR error for embedded image: {str(e)}")
        return ""

def extract_text_from_url(url):
    """Fetch and extract main text content from a URL."""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        # Remove scripts/styles
        for tag in soup(["script", "style", "nav", "footer", "header", "form", "noscript"]):
            tag.decompose()
        # Get visible text
        text = ' '.join(soup.stripped_strings)
        return text
    except Exception as e:
        logging.warning(f"Could not fetch or parse URL {url}: {str(e)}")
        return ""

def extract_urls_from_text(text):
    """Extract URLs from plain text using regex."""
    url_pattern = r'(https?://[^\s\)\]\}\"\']+)'
    return re.findall(url_pattern, text)

def extract_urls_from_docx(doc):
    """Extract URLs from a python-docx Document object."""
    urls = []
    for para in doc.paragraphs:
        if para.hyperlinks:
            for link in para.hyperlinks:
                urls.append(link.target)
        urls += extract_urls_from_text(para.text)
    return urls

def extract_urls_from_pptx(prs):
    """Extract URLs from a python-pptx Presentation object."""
    urls = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                urls += extract_urls_from_text(shape.text)
            if hasattr(shape, "hyperlink") and shape.hyperlink.address:
                urls.append(shape.hyperlink.address)
    return urls

# --- Update process_document to crawl and add referenced URLs ---
def process_document(file_path):
    """Process PDF, Word, PowerPoint, or image file and create vector database, including referenced URLs."""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        text = ""
        url_texts = []
        urls = []
        if ext == '.pdf':
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text() or ""
                    # Fallback to OCR if no text is found
                    if not page_text.strip() and '/XObject' in page.get('/Resources', {}):
                        xObject = page['/Resources']['/XObject'].get_object()
                        for obj in xObject:
                            if xObject[obj]['/Subtype'] == '/Image':
                                data = xObject[obj]._data
                                ocr_text = extract_text_from_image_bytes(data)
                                page_text += ocr_text
                    elif not page_text.strip():
                        pass
                    text += page_text
                    urls += extract_urls_from_text(page_text)
        elif ext == '.docx':
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
            # OCR for images in docx
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_bytes = rel.target_part.blob
                    ocr_text = extract_text_from_image_bytes(image_bytes)
                    text += ocr_text + "\n"
            # Extract URLs
            urls += extract_urls_from_docx(doc)
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    # OCR for images in pptx
                    if shape.shape_type == 13:  # PICTURE
                        image = shape.image
                        image_bytes = image.blob
                        ocr_text = extract_text_from_image_bytes(image_bytes)
                        text += ocr_text + "\n"
            # Extract URLs
            urls += extract_urls_from_pptx(prs)
        elif ext in ['.png', '.jpg', '.jpeg']:
            text = extract_text_from_image_file(file_path)
        else:
            raise ValueError("Unsupported file type.")

        # Fetch and add referenced URL content
        for url in set(urls):
            url_content = extract_text_from_url(url)
            if url_content:
                url_texts.append(f"[Referenced from {url}]\n{url_content}")

        # Combine all text
        all_text = text + "\n".join(url_texts)

        if not all_text.strip():
            raise ValueError("No text could be extracted from the document or its references.")

        # Split text
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        docs = splitter.create_documents([all_text])

        # Get API key
        cohere_api_key = os.environ.get("COHERE_API_KEY")
        if not cohere_api_key:
            raise ValueError("COHERE_API_KEY not found in environment variables")

        # Create embeddings and vector store
        embeddings = CohereEmbeddings(
            model="embed-english-v3.0",
            cohere_api_key=cohere_api_key
        )
        vectordb = Chroma.from_documents(docs, embeddings)

        return vectordb, len(docs)
    except Exception as e:
        logging.error(f"Error processing document: {str(e)}")
        raise

@app.route('/')
def index():
    """Main page"""
    # Initialize session messages if not exists
    if 'messages' not in session:
        session['messages'] = []
    
    # Check if API key is available
    cohere_api_key = os.environ.get("COHERE_API_KEY")
    
    return render_template('index.html', 
                         has_api_key=bool(cohere_api_key),
                         messages=session.get('messages', []))

@app.route('/upload', methods=['POST'])
def upload_file():
    """Handle PDF file upload"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file selected'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            file.save(file_path)
            
            # Process document (PDF, Word, or PowerPoint)
            vectordb, chunk_count = process_document(file_path)
            
            # Store vector database in session (in production, use a proper database)
            session['pdf_processed'] = True
            session['pdf_filename'] = filename
            session['chunk_count'] = chunk_count
            
            # Store vectordb reference (in production, use proper persistence)
            app.vectordb = vectordb
            
            # Clean up uploaded file
            os.remove(file_path)
            
            return jsonify({
                'success': True, 
                'message': f'PDF processed successfully! Created {chunk_count} text chunks.',
                'filename': filename
            })
        
        else:
            return jsonify({'error': 'Invalid file type. Please upload a PDF, Word, PowerPoint, or image file (pdf, docx, pptx, png, jpg, jpeg).'}), 400
    
    except Exception as e:
        logging.error(f"Upload error: {str(e)}")
        return jsonify({'error': f'Error processing PDF: {str(e)}'}), 500

@app.route('/ask', methods=['POST'])
def ask_question():
    """Handle question asking"""
    try:
        data = request.get_json()
        question = data.get('question', '').strip()
        
        if not question:
            return jsonify({'error': 'Please enter a question'}), 400
        
        if not session.get('pdf_processed'):
            return jsonify({'error': 'Please upload a document (PDF, Word, or PowerPoint) first'}), 400
        
        if not hasattr(app, 'vectordb'):
            return jsonify({'error': 'PDF processing session expired. Please upload the file again.'}), 400
        
        # Get API key
        cohere_api_key = os.environ.get("COHERE_API_KEY")
        if not cohere_api_key:
            return jsonify({'error': 'COHERE API key not configured'}), 500
        
        # Create QA chain
        retriever = app.vectordb.as_retriever()
        llm = init_chat_model("command-r-plus", model_provider="cohere", api_key=cohere_api_key)
        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=retriever,
            return_source_documents=False,
        )
        
        # Get answer
        response = qa_chain.invoke({"query": question})["result"]
        
        # Update session messages
        if 'messages' not in session:
            session['messages'] = []
        
        session['messages'].append({"role": "user", "content": question})
        session['messages'].append({"role": "assistant", "content": response})
        
        # Keep only last 20 messages to prevent session from growing too large
        if len(session['messages']) > 20:
            session['messages'] = session['messages'][-20:]
        
        session.modified = True
        
        return jsonify({
            'success': True,
            'answer': response,
            'question': question
        })
    
    except Exception as e:
        logging.error(f"Question answering error: {str(e)}")
        return jsonify({'error': f'Error getting answer: {str(e)}'}), 500

@app.route('/clear', methods=['POST'])
def clear_chat():
    """Clear chat history"""
    session['messages'] = []
    session.modified = True
    return jsonify({'success': True})

@app.route('/reset', methods=['POST'])
def reset_session():
    """Reset entire session"""
    session.clear()
    if hasattr(app, 'vectordb'):
        delattr(app, 'vectordb')
    return jsonify({'success': True})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
